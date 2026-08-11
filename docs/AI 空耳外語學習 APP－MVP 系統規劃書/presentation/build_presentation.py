# -*- coding: utf-8 -*-
"""Build Soraeru MVP PowerPoint + HTML deck with Stitch screenshots."""

from __future__ import annotations

import math
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
SCREENS = ASSETS / "screens"
DIAG = ASSETS / "diagrams"
OUT_PPTX = ROOT / "Soraeru_MVP_簡報.pptx"
OUT_HTML = ROOT / "Soraeru_MVP_簡報.html"

# Brand colors from Stitch DESIGN
TEAL = RGBColor(0x00, 0x4D, 0x64)
TEAL_MID = RGBColor(0x00, 0x66, 0x84)
TEAL_SOFT = RGBColor(0x87, 0xD0, 0xF2)
SLATE = RGBColor(0x4D, 0x61, 0x6C)
INK = RGBColor(0x18, 0x1C, 0x1E)
MUTED = RGBColor(0x3F, 0x48, 0x4D)
BG = RGBColor(0xF5, 0xFA, 0xFD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0x8B, 0x50, 0x00)
INFO = RGBColor(0x00, 0x61, 0xA4)
CARD = RGBColor(0xEC, 0xEE, 0xF1)
LINE = RGBColor(0xBF, 0xC8, 0xCD)

W, H = Inches(13.333), Inches(7.5)  # 16:9


def find_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        r"C:\Windows\Fonts\msjhbd.ttc" if bold else r"C:\Windows\Fonts\msjh.ttc",
        r"C:\Windows\Fonts\msyhbd.ttc" if bold else r"C:\Windows\Fonts\msyh.ttc",
        r"C:\Windows\Fonts\arial.ttf",
    ]
    for path in candidates:
        p = Path(path)
        if p.exists():
            try:
                return ImageFont.truetype(str(p), size=size, index=0)
            except OSError:
                continue
    return ImageFont.load_default()


def rounded_rect(draw, xy, radius, fill, outline=None, width=1):
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def make_diagrams():
    DIAG.mkdir(parents=True, exist_ok=True)
    make_value_loop(DIAG / "value_loop.png")
    make_user_flow(DIAG / "user_flow.png")
    make_architecture(DIAG / "architecture.png")
    make_timeline(DIAG / "timeline.png")
    make_scope_visual(DIAG / "scope.png")


def make_value_loop(path: Path):
    img = Image.new("RGB", (1600, 700), "#F5FAFD")
    d = ImageDraw.Draw(img)
    title_f = find_font(36, True)
    body_f = find_font(26)
    small_f = find_font(22)
    d.text((60, 40), "核心閉環：發音 → 空耳 → 記憶", font=title_f, fill="#004D64")

    steps = [
        ("1", "輸入／OCR", "手動或拍照取字", "#004D64"),
        ("2", "多語偵測", "AI 判斷來源語言", "#006684"),
        ("3", "空耳候選", "2～3 個華語近似音", "#4D616C"),
        ("4", "選存單字卡", "綁定帳號收藏", "#8B5000"),
    ]
    x0, y0, box_w, box_h, gap = 60, 160, 300, 280, 70
    for i, (num, title, desc, color) in enumerate(steps):
        x = x0 + i * (box_w + gap)
        rounded_rect(d, (x, y0, x + box_w, y0 + box_h), 28, "#FFFFFF", "#BFC8CD", 2)
        d.ellipse((x + 24, y0 + 28, x + 88, y0 + 92), fill=color)
        d.text((x + 46, y0 + 42), num, font=title_f, fill="#FFFFFF")
        d.text((x + 24, y0 + 130), title, font=title_f, fill="#181C1E")
        d.text((x + 24, y0 + 190), desc, font=body_f, fill="#3F484D")
        if i < len(steps) - 1:
            ax = x + box_w + 12
            ay = y0 + box_h // 2
            d.polygon([(ax, ay - 14), (ax + 40, ay), (ax, ay + 14)], fill="#87D0F2")
    d.text((60, 500), "正式發音可播放｜近似音僅供記憶｜原圖只在手機端 OCR", font=small_f, fill="#0061A4")
    img.save(path)


def make_user_flow(path: Path):
    img = Image.new("RGB", (1600, 720), "#F5FAFD")
    d = ImageDraw.Draw(img)
    title_f = find_font(34, True)
    body_f = find_font(22)
    d.text((50, 30), "使用者旅程", font=title_f, fill="#004D64")

    nodes = [
        [(180, 120, "Splash"), (480, 120, "登入／註冊"), (780, 120, "首次說明"), (1100, 120, "首頁")],
        [(260, 340, "輸入單字"), (620, 340, "拍照 OCR"), (980, 340, "分析中"), (1320, 340, "結果頁")],
        [(620, 560, "單字卡列表"), (980, 560, "詳情／播放"), (1320, 560, "設定／額度")],
    ]
    # draw nodes
    for row in nodes:
        for x, y, label in row:
            w, h = 220, 78
            rounded_rect(d, (x - w // 2, y - h // 2, x + w // 2, y + h // 2), 18, "#004D64")
            tw = d.textlength(label, font=body_f)
            d.text((x - tw / 2, y - 14), label, font=body_f, fill="#FFFFFF")

    def arrow(a, b):
        x1, y1 = a
        x2, y2 = b
        d.line((x1, y1, x2, y2), fill="#87D0F2", width=6)
        ang = math.atan2(y2 - y1, x2 - x1)
        for da in (2.6, -2.6):
            lx = x2 - 18 * math.cos(ang + da)
            ly = y2 - 18 * math.sin(ang + da)
            d.line((x2, y2, lx, ly), fill="#87D0F2", width=6)

    arrow((290, 120), (370, 120))
    arrow((590, 120), (670, 120))
    arrow((890, 120), (990, 120))
    arrow((1100, 160), (260, 300))
    arrow((1100, 160), (620, 300))
    arrow((370, 340), (870, 340))
    arrow((730, 340), (870, 340))
    arrow((1090, 340), (1210, 340))
    arrow((1320, 380), (980, 520))
    arrow((1320, 380), (620, 520))
    arrow((730, 560), (870, 560))
    img.save(path)


def make_architecture(path: Path):
    img = Image.new("RGB", (1600, 780), "#F5FAFD")
    d = ImageDraw.Draw(img)
    title_f = find_font(34, True)
    h_f = find_font(26, True)
    body_f = find_font(22)
    d.text((50, 28), "系統架構（MVP）", font=title_f, fill="#004D64")

    groups = [
        (50, 110, 480, 620, "Android 手機", "#D0E6F3", [
            "App UI（MAUI）",
            "裝置端 OCR",
            "系統 TTS",
            "本機快取",
        ]),
        (560, 110, 480, 620, "雲端", "#BEE9FF", [
            "ASP.NET Minimal API",
            "Auth（Google＋Email）",
            "Users／WordCards／Quota",
            "程序內快取",
        ]),
        (1070, 110, 480, 620, "外部服務", "#FFDCC0", [
            "LLM API",
            "Google IdP",
            "郵件（重設密碼）",
            "Google Play",
        ]),
    ]
    for x, y, w, h, title, fill, items in groups:
        rounded_rect(d, (x, y, x + w, y + h), 28, fill)
        d.text((x + 28, y + 24), title, font=h_f, fill="#004D64")
        for i, item in enumerate(items):
            iy = y + 100 + i * 100
            rounded_rect(d, (x + 28, iy, x + w - 28, iy + 72), 16, "#FFFFFF")
            d.text((x + 48, iy + 22), item, font=body_f, fill="#181C1E")

    # arrows between groups
    for y in (250, 450):
        d.polygon([(530, y - 12), (560, y), (530, y + 12)], fill="#004D64")
        d.polygon([(1040, y - 12), (1070, y), (1040, y + 12)], fill="#004D64")
    img.save(path)


def make_timeline(path: Path):
    img = Image.new("RGB", (1600, 700), "#F5FAFD")
    d = ImageDraw.Draw(img)
    title_f = find_font(34, True)
    body_f = find_font(20)
    small_f = find_font(18)
    d.text((50, 30), "8 週開發時程", font=title_f, fill="#004D64")
    weeks = [
        ("W1", "Auth＋登入註冊"),
        ("W2", "Analyze API＋Prompt"),
        ("W3", "主閉環 UI＋單字卡"),
        ("W4", "OCR＋TTS＋多語"),
        ("W5", "額度／設定／視覺"),
        ("W6-7", "封閉測試 12×14"),
        ("W8", "商店素材＋送審"),
    ]
    d.line((80, 280, 1520, 280), fill="#87D0F2", width=10)
    for i, (w, label) in enumerate(weeks):
        x = 120 + i * 200
        d.ellipse((x - 18, 262, x + 18, 298), fill="#004D64")
        rounded_rect(d, (x - 80, 140 if i % 2 == 0 else 340, x + 80, 220 if i % 2 == 0 else 420), 16, "#FFFFFF", "#BFC8CD", 2)
        ty = 155 if i % 2 == 0 else 355
        d.text((x - 28, ty), w, font=title_f, fill="#004D64")
        # wrap-ish
        d.text((x - 72, ty + 42), label[:8], font=body_f, fill="#181C1E")
        if len(label) > 8:
            d.text((x - 72, ty + 68), label[8:], font=body_f, fill="#181C1E")
    d.text((50, 520), "工時粗估 300～420 小時｜現金預備建議 2～4 萬 TWD", font=small_f, fill="#3F484D")
    d.text((50, 560), "定案：VS2026 ＋ .NET MAUI ＋ ASP.NET Core ＋ Stitch／AI Studio", font=small_f, fill="#006684")
    img.save(path)


def make_scope_visual(path: Path):
    img = Image.new("RGB", (1600, 780), "#F5FAFD")
    d = ImageDraw.Draw(img)
    title_f = find_font(34, True)
    h_f = find_font(28, True)
    body_f = find_font(24)
    d.text((50, 30), "MVP 範圍：做與不做", font=title_f, fill="#004D64")

    rounded_rect(d, (50, 110, 760, 720), 28, "#E8F7EE")
    rounded_rect(d, (840, 110, 1550, 720), 28, "#FDECEC")
    d.text((90, 140), "✓ 納入首版", font=h_f, fill="#0B6E3F")
    d.text((880, 140), "✕ 明確砍除", font=h_f, fill="#BA1A1A")

    yes = [
        "Google／Email 帳號與日額度",
        "手動輸入＋裝置端 OCR",
        "多語自動偵測空耳",
        "詞義／讀音／TTS／2～3 候選",
        "帳號綁定單字卡雲端同步",
        "Android／Google Play",
    ]
    no = [
        "多 Agent 編排",
        "iOS／App Store",
        "雲端 OCR／上傳原圖",
        "完整 SRS／間隔複習",
        "實際金流／Billing",
        "社群／排行榜／梗圖生圖",
    ]
    for i, t in enumerate(yes):
        d.text((90, 220 + i * 70), f"•  {t}", font=body_f, fill="#181C1E")
    for i, t in enumerate(no):
        d.text((880, 220 + i * 70), f"•  {t}", font=body_f, fill="#181C1E")
    img.save(path)


def set_run(run, size=18, bold=False, color=INK, font="Microsoft JhengHei"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.get_or_add_eastAsian() if hasattr(rPr, "get_or_add_eastAsian") else None
    # Explicit east asian font
    rFonts = rPr.find(qn("a:rFonts"))
    if rFonts is None:
        from lxml import etree

        rFonts = etree.SubElement(rPr, qn("a:rFonts"))
    rFonts.set("ascii", font)
    rFonts.set("hAnsi", font)
    rFonts.set("eastAsia", font)


def add_bg(slide, color=BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # send to back
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_accent_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()


def add_footer(slide, page: int, total: int):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(10), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Soraeru｜空耳學單字 MVP ｜ v1.1"
    set_run(run, 11, False, MUTED)
    num = slide.shapes.add_textbox(Inches(11.5), Inches(7.05), Inches(1.4), Inches(0.35))
    p2 = num.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = f"{page} / {total}"
    set_run(r2, 11, False, MUTED)


def add_title(slide, title: str, subtitle: str | None = None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run(r, 32, True, TEAL)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.55), Inches(0.9), Inches(12), Inches(0.4))
        sp = sub.text_frame.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        set_run(sr, 14, False, MUTED)


def add_text_box(slide, left, top, width, height, text, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_run(r, size, bold, color)
    return box


def add_card(slide, left, top, width, height, fill=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    shape.adjustments[0] = 0.08
    return shape


def add_phone(slide, img_path: Path, left, top, height=Inches(5.6)):
    # phone frame
    phone_w = height * 390 / 844
    frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left - Inches(0.06),
        top - Inches(0.06),
        phone_w + Inches(0.12),
        height + Inches(0.12),
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = INK
    frame.line.fill.background()
    frame.adjustments[0] = 0.1
    slide.shapes.add_picture(str(img_path), left, top, height=height)
    return phone_w


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_pptx():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    total = 17
    page = 0

    def next_slide():
        nonlocal page
        page += 1
        s = blank_slide(prs)
        add_bg(s)
        add_accent_bar(s)
        return s

    # 1 Cover
    s = blank_slide(prs)
    page = 1
    # gradient-like panels
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = TEAL
    bg.line.fill.background()
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.2), 0, Inches(6.2), H)
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(0x00, 0x3A, 0x4C)
    panel.line.fill.background()
    add_text_box(s, Inches(0.8), Inches(1.6), Inches(6), Inches(0.5), "Soraeru  ·  Google Play MVP", 16, False, TEAL_SOFT)
    add_text_box(s, Inches(0.8), Inches(2.2), Inches(6.2), Inches(1.2), "空耳學單字", 54, True, WHITE)
    add_text_box(s, Inches(0.8), Inches(3.4), Inches(6), Inches(0.8), "用發音，記住外語", 24, False, TEAL_SOFT)
    add_text_box(
        s,
        Inches(0.8),
        Inches(4.4),
        Inches(6),
        Inches(1.2),
        "外語發音 → 台灣華語空耳近似音\n加速記憶的小型語言學習 App",
        16,
        False,
        WHITE,
    )
    add_text_box(s, Inches(0.8), Inches(6.5), Inches(6), Inches(0.4), "規劃書 v1.1 ｜ 2026-08-06", 12, False, TEAL_SOFT)
    # phone on cover
    if (SCREENS / "L05_home.png").exists():
        add_phone(s, SCREENS / "L05_home.png", Inches(8.55), Inches(0.85), Inches(5.9))
    add_text_box(s, Inches(11.5), Inches(7.05), Inches(1.4), Inches(0.3), f"{page}/{total}", 11, False, TEAL_SOFT, PP_ALIGN.RIGHT)

    # 2 Agenda
    s = next_slide()
    add_title(s, "簡報大綱", "從產品定位到可開工交付")
    items = [
        ("01", "產品定位與核心價值"),
        ("02", "MVP 範圍與功能亮點"),
        ("03", "使用者旅程與畫面設計"),
        ("04", "系統架構與技術棧"),
        ("05", "時程、指標與開工清單"),
    ]
    for i, (num, title) in enumerate(items):
        y = Inches(1.55) + Inches(i * 0.95)
        add_card(s, Inches(0.7), y, Inches(11.8), Inches(0.8), WHITE)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), y + Inches(0.15), Inches(0.5), Inches(0.5))
        circ.fill.solid()
        circ.fill.fore_color.rgb = TEAL
        circ.line.fill.background()
        add_text_box(s, Inches(0.95), y + Inches(0.25), Inches(0.5), Inches(0.35), num[-1], 16, True, WHITE, PP_ALIGN.CENTER)
        add_text_box(s, Inches(1.7), y + Inches(0.22), Inches(10), Inches(0.45), title, 22, True, INK)
    add_footer(s, page, total)

    # 3 Product definition
    s = next_slide()
    add_title(s, "一句話產品定義", "範圍凍結後的開工規格")
    add_card(s, Inches(0.7), Inches(1.5), Inches(7.4), Inches(4.8), WHITE)
    add_text_box(
        s,
        Inches(1.0),
        Inches(1.8),
        Inches(6.8),
        Inches(4.2),
        "使用者登入後，以手動輸入或拍照／相簿 OCR 選字，"
        "經單一 AI 分析 API 自動判斷來源語言，取得詞義、正式讀音與台灣華語近似音候選；"
        "選一個候選後存入與帳號綁定的單字卡。\n\n"
        "重點不是「只學英美日」，而是：\n"
        "任何外語發音 → 華語可記的空耳 → 加速記憶。",
        18,
        False,
        INK,
    )
    add_card(s, Inches(8.4), Inches(1.5), Inches(4.3), Inches(1.4), RGBColor(0xD0, 0xE6, 0xF3))
    add_text_box(s, Inches(8.6), Inches(1.7), Inches(3.9), Inches(1.0), "平台\nAndroid ＋ Google Play", 18, True, TEAL)
    add_card(s, Inches(8.4), Inches(3.15), Inches(4.3), Inches(1.4), RGBColor(0xBE, 0xE9, 0xFF))
    add_text_box(s, Inches(8.6), Inches(3.35), Inches(3.9), Inches(1.0), "語言\n多語自動偵測，不鎖白名單", 18, True, TEAL)
    add_card(s, Inches(8.4), Inches(4.8), Inches(4.3), Inches(1.5), RGBColor(0xFF, 0xDC, 0xC0))
    add_text_box(s, Inches(8.6), Inches(5.0), Inches(3.9), Inches(1.1), "帳戶\nGoogle＋Email／日額度", 18, True, TEAL)
    add_footer(s, page, total)

    # 4 Value loop
    s = next_slide()
    add_title(s, "核心價值", "系統產初稿，使用者只判斷好不好記")
    s.shapes.add_picture(str(DIAG / "value_loop.png"), Inches(0.5), Inches(1.35), width=Inches(12.3))
    add_footer(s, page, total)

    # 5 Scope
    s = next_slide()
    add_title(s, "MVP 範圍", "先求可上架閉環，砍掉不可控成本")
    s.shapes.add_picture(str(DIAG / "scope.png"), Inches(0.45), Inches(1.2), width=Inches(12.4))
    add_footer(s, page, total)

    # 6 Features
    s = next_slide()
    add_title(s, "功能亮點", "F01–F20 中的體驗主軸")
    feats = [
        ("⌨", "手動輸入", "≤50 字短語\n預設自動偵測語言"),
        ("📷", "裝置端 OCR", "拍照／相簿選字\n原圖絕不上傳"),
        ("🌐", "多語空耳", "詞義＋讀音＋\n2～3 近似音候選"),
        ("▶", "正式發音", "系統 TTS 播放\n近似音不播"),
        ("📚", "單字卡雲端", "搜尋／語言篩選\n帳號歸屬查重"),
        ("⚡", "額度控費", "每日 AI 上限\n重產≤3／快取命中"),
    ]
    for i, (icon, title, desc) in enumerate(feats):
        col, row = i % 3, i // 3
        left = Inches(0.7) + Inches(col * 4.1)
        top = Inches(1.5) + Inches(row * 2.5)
        add_card(s, left, top, Inches(3.85), Inches(2.2), WHITE)
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.25), top + Inches(0.25), Inches(0.7), Inches(0.55))
        badge.fill.solid()
        badge.fill.fore_color.rgb = TEAL
        badge.line.fill.background()
        add_text_box(s, left + Inches(0.25), top + Inches(0.32), Inches(0.7), Inches(0.4), icon, 16, True, WHITE, PP_ALIGN.CENTER)
        add_text_box(s, left + Inches(0.25), top + Inches(0.95), Inches(3.3), Inches(0.4), title, 20, True, TEAL)
        add_text_box(s, left + Inches(0.25), top + Inches(1.4), Inches(3.3), Inches(0.7), desc, 14, False, MUTED)
    add_footer(s, page, total)

    # 7 User flow
    s = next_slide()
    add_title(s, "使用者旅程", "登入 → 輸入／OCR → 空耳 → 收藏")
    s.shapes.add_picture(str(DIAG / "user_flow.png"), Inches(0.4), Inches(1.2), width=Inches(12.5))
    add_footer(s, page, total)

    # 8 UI auth
    s = next_slide()
    add_title(s, "關鍵畫面｜進入體驗", "Splash → 登入 → 首次說明（L00／L01／L04）")
    phones = [
        (SCREENS / "L00_splash.png", "品牌啟動"),
        (SCREENS / "L01_login.png", "Google／Email 登入"),
        (SCREENS / "L04_onboarding.png", "首次使用說明"),
    ]
    x = Inches(0.7)
    for path, label in phones:
        if path.exists():
            pw = add_phone(s, path, x, Inches(1.35), Inches(5.2))
            add_text_box(s, x, Inches(6.65), pw, Inches(0.3), label, 13, True, TEAL, PP_ALIGN.CENTER)
            x += pw + Inches(0.55)
    add_footer(s, page, total)

    # 9 UI showcase home+input+result
    s = next_slide()
    add_title(s, "關鍵畫面｜主閉環", "Stitch 高保真 UI（L05／L06／L10）")
    phones = [
        (SCREENS / "L05_home.png", "首頁入口"),
        (SCREENS / "L06_input.png", "單字輸入"),
        (SCREENS / "L10_result.png", "分析結果"),
    ]
    x = Inches(0.7)
    for path, label in phones:
        if path.exists():
            pw = add_phone(s, path, x, Inches(1.35), Inches(5.2))
            add_text_box(s, x, Inches(6.65), pw, Inches(0.3), label, 13, True, TEAL, PP_ALIGN.CENTER)
            x += pw + Inches(0.55)
    add_footer(s, page, total)

    # 9 OCR path
    s = next_slide()
    add_title(s, "關鍵畫面｜圖片取字", "隱私優先：只在手機辨識，後端只收文字")
    phones = [
        (SCREENS / "L07_image.png", "選圖／拍照"),
        (SCREENS / "L08_ocr.png", "OCR 選字"),
        (SCREENS / "L09_analyzing.png", "分析中"),
    ]
    x = Inches(0.7)
    for path, label in phones:
        if path.exists():
            pw = add_phone(s, path, x, Inches(1.35), Inches(5.2))
            add_text_box(s, x, Inches(6.65), pw, Inches(0.3), label, 13, True, TEAL, PP_ALIGN.CENTER)
            x += pw + Inches(0.55)
    add_footer(s, page, total)

    # 10 notebook/settings
    s = next_slide()
    add_title(s, "關鍵畫面｜帳號與收藏", "單字卡、詳情、設定與額度預留")
    phones = [
        (SCREENS / "L11_notebook.png", "我的單字卡"),
        (SCREENS / "L12_detail.png", "單字卡詳情"),
        (SCREENS / "L13_settings.png", "設定／帳號"),
    ]
    x = Inches(0.7)
    for path, label in phones:
        if path.exists():
            pw = add_phone(s, path, x, Inches(1.35), Inches(5.2))
            add_text_box(s, x, Inches(6.65), pw, Inches(0.3), label, 13, True, TEAL, PP_ALIGN.CENTER)
            x += pw + Inches(0.55)
    add_footer(s, page, total)

    # 11 Design system
    s = next_slide()
    add_title(s, "設計系統｜Soraeru", "Pragmatic & Supportive｜Deep Teal 專業學習工具")
    # color swatches
    swatches = [
        (TEAL, "Primary\n#004D64"),
        (TEAL_MID, "Container\n#006684"),
        (TEAL_SOFT, "Soft\n#87D0F2"),
        (SLATE, "Secondary\n#4D616C"),
        (AMBER, "Warning\n#8B5000"),
        (INFO, "Info\n#0061A4"),
    ]
    for i, (color, label) in enumerate(swatches):
        left = Inches(0.7) + Inches(i * 2.05)
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(1.9), Inches(1.5))
        sh.fill.solid()
        sh.fill.fore_color.rgb = color
        sh.line.fill.background()
        add_text_box(s, left, Inches(3.15), Inches(1.9), Inches(0.7), label, 12, True, INK, PP_ALIGN.CENTER)

    principles = [
        ("單一任務", "每頁最多一個實心 Primary CTA"),
        ("安全優先", "空耳警示與 OCR 隱私視為一等 UI"),
        ("層次深度", "用 surface 層級區分內容，不用厚重陰影"),
        ("雙語可讀", "Hanken Grotesk + Noto Sans + JetBrains Mono"),
    ]
    for i, (t, d_) in enumerate(principles):
        col, row = i % 2, i // 2
        left = Inches(0.7) + Inches(col * 6.2)
        top = Inches(4.1) + Inches(row * 1.2)
        add_card(s, left, top, Inches(5.95), Inches(1.05), WHITE)
        add_text_box(s, left + Inches(0.25), top + Inches(0.15), Inches(5.4), Inches(0.35), t, 16, True, TEAL)
        add_text_box(s, left + Inches(0.25), top + Inches(0.5), Inches(5.4), Inches(0.4), d_, 13, False, MUTED)
    add_footer(s, page, total)

    # 12 Architecture
    s = next_slide()
    add_title(s, "系統架構", "手機端 OCR／TTS ＋ 雲端單一分析 API")
    s.shapes.add_picture(str(DIAG / "architecture.png"), Inches(0.4), Inches(1.15), width=Inches(12.5))
    add_footer(s, page, total)

    # 13 Tech stack
    s = next_slide()
    add_title(s, "技術棧與工具鏈", "以 Visual Studio 2026 為主力 IDE")
    stacks = [
        ("App", ".NET MAUI\nAndroid"),
        ("API", "ASP.NET Core\nMinimal API"),
        ("Auth", "Google Sign-In\n＋ Email／密碼"),
        ("OCR", "裝置端\nML Kit 等同"),
        ("LLM", "AI Studio\n選定文字模型"),
        ("UI", "Google Stitch\nDESIGN.md"),
    ]
    for i, (k, v) in enumerate(stacks):
        col, row = i % 3, i // 3
        left = Inches(0.7) + Inches(col * 4.1)
        top = Inches(1.5) + Inches(row * 2.35)
        add_card(s, left, top, Inches(3.9), Inches(2.1), WHITE)
        add_text_box(s, left + Inches(0.3), top + Inches(0.35), Inches(3.3), Inches(0.4), k, 14, False, MUTED)
        add_text_box(s, left + Inches(0.3), top + Inches(0.85), Inches(3.3), Inches(1.0), v, 22, True, TEAL)
    add_footer(s, page, total)

    # 14 Timeline
    s = next_slide()
    add_title(s, "開發時程", "目標 8 週上架軌道")
    s.shapes.add_picture(str(DIAG / "timeline.png"), Inches(0.45), Inches(1.25), width=Inches(12.4))
    add_footer(s, page, total)

    # 15 Metrics & privacy
    s = next_slide()
    add_title(s, "成功指標與隱私原則", "可量測、可上架、可信賴")
    metrics = [
        ("登入成功率", "≥ 98%"),
        ("生成成功率", "≥ 95%"),
        ("收藏率", "≥ 40%"),
        ("P95 延遲", "< 8 秒"),
        ("Crash-free", "≥ 99%"),
    ]
    for i, (k, v) in enumerate(metrics):
        left = Inches(0.55) + Inches(i * 2.5)
        add_card(s, left, Inches(1.5), Inches(2.35), Inches(1.8), WHITE)
        add_text_box(s, left + Inches(0.1), Inches(1.7), Inches(2.15), Inches(0.4), k, 13, False, MUTED, PP_ALIGN.CENTER)
        add_text_box(s, left + Inches(0.1), Inches(2.25), Inches(2.15), Inches(0.7), v, 26, True, TEAL, PP_ALIGN.CENTER)

    add_card(s, Inches(0.7), Inches(3.7), Inches(5.9), Inches(2.8), RGBColor(0xD0, 0xE6, 0xF3))
    add_text_box(s, Inches(0.95), Inches(3.95), Inches(5.4), Inches(0.4), "隱私與資安", 18, True, TEAL)
    add_text_box(
        s,
        Inches(0.95),
        Inches(4.45),
        Inches(5.4),
        Inches(1.8),
        "• 原圖不離開手機\n• 僅選定文字送 API／LLM\n• LLM Key 不上 App\n• JWT／密碼雜湊／最小權限",
        15,
        False,
        INK,
    )
    add_card(s, Inches(6.9), Inches(3.7), Inches(5.8), Inches(2.8), RGBColor(0xFF, 0xDC, 0xC0))
    add_text_box(s, Inches(7.15), Inches(3.95), Inches(5.3), Inches(0.4), "風險對策", 18, True, AMBER)
    add_text_box(
        s,
        Inches(7.15),
        Inches(4.45),
        Inches(5.3),
        Inches(1.8),
        "• 多語品質不均 → 聲明＋重產＋覆寫語言\n• OCR 弱語系 → 手動輸入一等公民\n• AI 費用 → 額度＋快取＋熔斷\n• 金流過早 → Phase 2 再接 Billing",
        15,
        False,
        INK,
    )
    add_footer(s, page, total)

    # 16 Closing
    s = next_slide()
    add_title(s, "開工閉環與下一步", "規格已合流，可直接動手")
    add_card(s, Inches(0.7), Inches(1.5), Inches(12), Inches(1.5), TEAL)
    add_text_box(
        s,
        Inches(1.0),
        Inches(1.85),
        Inches(11.4),
        Inches(0.9),
        "登入 → 輸入或 OCR 選字 → 多語偵測與空耳 → 選候選 → 存入帳號單字卡",
        22,
        True,
        WHITE,
        PP_ALIGN.CENTER,
    )
    steps = [
        "1. Auth＋L01～L03",
        "2. Users／Quota／JWT",
        "3. 多語 Prompt 定稿",
        "4. Analyze＋結果頁",
        "5. Notebook＋OCR＋TTS",
        "6. 封閉測試＋送審",
    ]
    for i, t in enumerate(steps):
        col, row = i % 3, i // 3
        left = Inches(0.7) + Inches(col * 4.1)
        top = Inches(3.4) + Inches(row * 1.35)
        add_card(s, left, top, Inches(3.9), Inches(1.15), WHITE)
        add_text_box(s, left + Inches(0.25), top + Inches(0.35), Inches(3.4), Inches(0.5), t, 16, True, TEAL)

    add_footer(s, page, total)

    prs.save(OUT_PPTX)
    print(f"PPTX -> {OUT_PPTX}")


def build_html():
    # relative paths from presentation folder
    def img(name: str) -> str:
        return f"assets/screens/{name}"

    html = f"""<!DOCTYPE html>
<html lang="zh-TW">
<head>
<meta charset="utf-8"/>
<meta name="viewport" content="width=device-width, initial-scale=1"/>
<title>Soraeru｜空耳學單字 MVP 簡報</title>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Hanken+Grotesk:wght@500;600;700&family=Noto+Sans+TC:wght@400;500;700&display=swap" rel="stylesheet">
<style>
  :root {{
    --teal:#004d64; --teal-mid:#006684; --soft:#87d0f2; --bg:#f5fafd;
    --ink:#181c1e; --muted:#3f484d; --card:#fff; --line:#bfc8cd;
    --amber:#8b5000; --info:#0061a4;
  }}
  * {{ box-sizing:border-box; margin:0; padding:0; }}
  html,body {{ height:100%; background:#0b1c24; color:var(--ink);
    font-family:"Noto Sans TC","Hanken Grotesk",sans-serif; overflow:hidden; }}
  .deck {{ height:100%; width:100%; position:relative; }}
  .slide {{
    position:absolute; inset:0; display:none; padding:48px 56px 64px;
    background:
      radial-gradient(1200px 600px at 10% -10%, rgba(135,208,242,.35), transparent 60%),
      radial-gradient(900px 500px at 100% 0%, rgba(0,102,132,.18), transparent 50%),
      linear-gradient(180deg, #f7fbfd 0%, #eef5f8 100%);
    animation: fade .35s ease;
  }}
  .slide.active {{ display:flex; flex-direction:column; }}
  .slide.cover {{
    background: linear-gradient(135deg, #004d64 0%, #003240 55%, #001f2a 100%);
    color:#fff; flex-direction:row; gap:40px; align-items:center;
  }}
  @keyframes fade {{ from {{ opacity:0; transform:translateY(8px);}} to {{opacity:1; transform:none;}} }}
  h1 {{ font-family:"Hanken Grotesk",sans-serif; font-size:42px; color:var(--teal); letter-spacing:-.02em; }}
  .cover h1 {{ color:#fff; font-size:64px; }}
  .sub {{ color:var(--muted); margin-top:8px; font-size:18px; }}
  .cover .sub {{ color:#a2e1ff; }}
  .eyebrow {{ color:var(--teal-mid); font-weight:700; letter-spacing:.08em; text-transform:uppercase; font-size:13px; }}
  .cover .eyebrow {{ color:#87d0f2; }}
  .grid {{ display:grid; gap:18px; margin-top:28px; }}
  .grid.2 {{ grid-template-columns:1.2fr .8fr; }}
  .grid.3 {{ grid-template-columns:repeat(3,1fr); }}
  .grid.6 {{ grid-template-columns:repeat(3,1fr); }}
  .card {{
    background:rgba(255,255,255,.9); border:1px solid var(--line); border-radius:18px;
    padding:22px; box-shadow:0 8px 24px rgba(0,77,100,.06); backdrop-filter:blur(8px);
  }}
  .card h3 {{ font-size:20px; color:var(--teal); margin-bottom:8px; }}
  .card p, .card li {{ color:var(--muted); line-height:1.6; font-size:15px; }}
  .card ul {{ padding-left:18px; }}
  .phones {{ display:flex; gap:28px; justify-content:center; align-items:flex-end; margin-top:18px; flex:1; }}
  .phone {{
    width:min(240px, 22vw); aspect-ratio:390/844; border-radius:28px; overflow:hidden;
    border:8px solid #181c1e; box-shadow:0 20px 50px rgba(0,0,0,.28); background:#fff;
    transform:translateY(0); transition:transform .3s ease;
  }}
  .phone:hover {{ transform:translateY(-8px); }}
  .phone img {{ width:100%; height:100%; object-fit:cover; display:block; }}
  .phone figcaption {{ text-align:center; margin-top:10px; color:var(--teal); font-weight:700; font-size:14px; }}
  .figure {{ width:100%; border-radius:16px; border:1px solid var(--line); background:#fff; margin-top:16px; }}
  .swatch {{ height:92px; border-radius:14px; margin-bottom:8px; }}
  .kpi {{ text-align:center; }}
  .kpi .v {{ font-size:34px; font-weight:800; color:var(--teal); font-family:"Hanken Grotesk",sans-serif; }}
  .nav {{
    position:fixed; left:0; right:0; bottom:0; height:48px; display:flex; align-items:center;
    justify-content:space-between; padding:0 24px; color:#cfe7f2; font-size:13px;
    background:linear-gradient(180deg, transparent, rgba(0,20,28,.55)); z-index:20;
  }}
  .hint {{ opacity:.8; }}
  .progress {{ position:fixed; top:0; left:0; height:4px; background:#87d0f2; z-index:30; transition:width .25s; }}
  .pill {{ display:inline-flex; align-items:center; gap:8px; background:#d0e6f3; color:var(--teal);
    border-radius:999px; padding:8px 14px; font-weight:700; font-size:14px; }}
  .big-quote {{
    margin-top:24px; padding:28px; border-radius:20px; background:var(--teal); color:#fff;
    font-size:24px; font-weight:700; text-align:center; line-height:1.5;
  }}
  .yes {{ background:#e8f7ee; }} .no {{ background:#fdecec; }}
  .agenda-item {{ display:flex; gap:16px; align-items:center; padding:18px 20px; }}
  .num {{ width:42px; height:42px; border-radius:50%; background:var(--teal); color:#fff;
    display:grid; place-items:center; font-weight:800; }}
  .cover-left {{ flex:1; padding-left:20px; }}
  .cover-right {{ flex:0 0 300px; }}
  .tagrow {{ display:flex; flex-wrap:wrap; gap:10px; margin-top:22px; }}
  .tag {{ background:rgba(255,255,255,.12); border:1px solid rgba(162,225,255,.35); color:#bee9ff;
    border-radius:999px; padding:8px 14px; font-size:13px; }}
</style>
</head>
<body>
<div class="progress" id="progress"></div>
<div class="deck" id="deck">

<section class="slide cover active">
  <div class="cover-left">
    <div class="eyebrow">SORAERU · GOOGLE PLAY MVP</div>
    <h1>空耳學單字</h1>
    <p class="sub">用發音，記住外語</p>
    <p style="margin-top:22px;max-width:560px;line-height:1.7;color:#d7eef8;font-size:18px;">
      外語發音 → 台灣華語空耳近似音 → 加速記憶。<br/>
      依 Cursor 合流規劃書與 Stitch 高保真 UI 打造可上架 MVP。
    </p>
    <div class="tagrow">
      <span class="tag">多語自動偵測</span>
      <span class="tag">裝置端 OCR</span>
      <span class="tag">帳號＋日額度</span>
      <span class="tag">.NET MAUI</span>
    </div>
  </div>
  <div class="cover-right">
    <figure class="phone" style="width:280px;margin:0 auto;"><img src="{img('L05_home.png')}" alt="home"/></figure>
  </div>
</section>

<section class="slide">
  <div class="eyebrow">AGENDA</div>
  <h1>簡報大綱</h1>
  <p class="sub">從產品定位到可開工交付</p>
  <div class="grid" style="margin-top:28px; gap:14px;">
    <div class="card agenda-item"><div class="num">1</div><div><h3>產品定位與核心價值</h3></div></div>
    <div class="card agenda-item"><div class="num">2</div><div><h3>MVP 範圍與功能亮點</h3></div></div>
    <div class="card agenda-item"><div class="num">3</div><div><h3>使用者旅程與畫面設計</h3></div></div>
    <div class="card agenda-item"><div class="num">4</div><div><h3>系統架構與技術棧</h3></div></div>
    <div class="card agenda-item"><div class="num">5</div><div><h3>時程、指標與開工清單</h3></div></div>
  </div>
</section>

<section class="slide">
  <div class="eyebrow">PRODUCT</div>
  <h1>一句話產品定義</h1>
  <div class="grid 2">
    <div class="card">
      <p style="font-size:18px;color:var(--ink);line-height:1.75;">
        使用者登入後，以<strong>手動輸入</strong>或<strong>拍照／相簿 OCR 選字</strong>，經<strong>單一 AI 分析 API</strong>自動判斷來源語言，
        取得詞義、正式讀音與台灣華語近似音候選；選一個候選後存入與帳號綁定的單字卡。
      </p>
      <p style="margin-top:18px;">核心不是「只學英美日」，而是任何外語的發音 → 華語可記的空耳 → 加速記憶。</p>
    </div>
    <div class="grid" style="gap:14px;">
      <div class="card" style="background:#d0e6f3;"><h3>平台</h3><p>Android ＋ Google Play</p></div>
      <div class="card" style="background:#bee9ff;"><h3>語言</h3><p>多語自動偵測，不鎖白名單</p></div>
      <div class="card" style="background:#ffdcc0;"><h3>帳戶</h3><p>Google＋Email／每日 AI 額度</p></div>
    </div>
  </div>
</section>

<section class="slide">
  <div class="eyebrow">VALUE</div>
  <h1>核心價值閉環</h1>
  <p class="sub">系統產初稿，使用者只判斷像不像、好不好記、要不要存</p>
  <img class="figure" src="assets/diagrams/value_loop.png" alt="value loop"/>
</section>

<section class="slide">
  <div class="eyebrow">SCOPE</div>
  <h1>MVP 範圍：做與不做</h1>
  <img class="figure" src="assets/diagrams/scope.png" alt="scope"/>
</section>

<section class="slide">
  <div class="eyebrow">FEATURES</div>
  <h1>功能亮點</h1>
  <div class="grid 6" style="margin-top:24px;">
    <div class="card"><h3>⌨ 手動輸入</h3><p>≤50 字短語，預設自動偵測語言，可手動覆寫。</p></div>
    <div class="card"><h3>📷 裝置端 OCR</h3><p>拍照／相簿選字，原圖不離開手機。</p></div>
    <div class="card"><h3>🌐 多語空耳</h3><p>詞義＋正式讀音＋2～3 個華語近似音候選。</p></div>
    <div class="card"><h3>▶ 正式發音</h3><p>系統 TTS 播放原文；近似音僅供記憶。</p></div>
    <div class="card"><h3>📚 單字卡雲端</h3><p>帳號歸屬、搜尋、語言篩選、查重。</p></div>
    <div class="card"><h3>⚡ 額度控費</h3><p>每日上限、同字重產≤3、快取命中不打 LLM。</p></div>
  </div>
</section>

<section class="slide">
  <div class="eyebrow">JOURNEY</div>
  <h1>使用者旅程</h1>
  <img class="figure" src="assets/diagrams/user_flow.png" alt="user flow"/>
</section>

<section class="slide">
  <div class="eyebrow">UI · AUTH</div>
  <h1>關鍵畫面｜進入體驗</h1>
  <p class="sub">Splash → 登入 → 首次說明</p>
  <div class="phones">
    <figure><div class="phone"><img src="{img('L00_splash.png')}"/></div><figcaption>品牌啟動</figcaption></figure>
    <figure><div class="phone"><img src="{img('L01_login.png')}"/></div><figcaption>Google／Email 登入</figcaption></figure>
    <figure><div class="phone"><img src="{img('L04_onboarding.png')}"/></div><figcaption>首次使用說明</figcaption></figure>
  </div>
</section>

<section class="slide">
  <div class="eyebrow">UI · CORE LOOP</div>
  <h1>關鍵畫面｜主閉環</h1>
  <p class="sub">Stitch 高保真：首頁 → 輸入 → 分析結果</p>
  <div class="phones">
    <figure><div class="phone"><img src="{img('L05_home.png')}"/></div><figcaption>首頁入口</figcaption></figure>
    <figure><div class="phone"><img src="{img('L06_input.png')}"/></div><figcaption>單字輸入</figcaption></figure>
    <figure><div class="phone"><img src="{img('L10_result.png')}"/></div><figcaption>分析結果</figcaption></figure>
  </div>
</section>

<section class="slide">
  <div class="eyebrow">UI · OCR</div>
  <h1>關鍵畫面｜圖片取字</h1>
  <p class="sub">隱私優先：後端只收選定文字</p>
  <div class="phones">
    <figure><div class="phone"><img src="{img('L07_image.png')}"/></div><figcaption>選圖／拍照</figcaption></figure>
    <figure><div class="phone"><img src="{img('L08_ocr.png')}"/></div><figcaption>OCR 選字</figcaption></figure>
    <figure><div class="phone"><img src="{img('L09_analyzing.png')}"/></div><figcaption>分析中</figcaption></figure>
  </div>
</section>

<section class="slide">
  <div class="eyebrow">UI · ACCOUNT</div>
  <h1>關鍵畫面｜帳號與收藏</h1>
  <div class="phones">
    <figure><div class="phone"><img src="{img('L11_notebook.png')}"/></div><figcaption>我的單字卡</figcaption></figure>
    <figure><div class="phone"><img src="{img('L12_detail.png')}"/></div><figcaption>單字卡詳情</figcaption></figure>
    <figure><div class="phone"><img src="{img('L13_settings.png')}"/></div><figcaption>設定／帳號</figcaption></figure>
  </div>
</section>

<section class="slide">
  <div class="eyebrow">DESIGN SYSTEM</div>
  <h1>設計系統｜Deep Teal</h1>
  <p class="sub">Pragmatic & Supportive · Functional Minimalism</p>
  <div class="grid" style="grid-template-columns:repeat(6,1fr); margin-top:20px;">
    <div><div class="swatch" style="background:#004d64"></div><p>Primary</p></div>
    <div><div class="swatch" style="background:#006684"></div><p>Container</p></div>
    <div><div class="swatch" style="background:#87d0f2"></div><p>Soft</p></div>
    <div><div class="swatch" style="background:#4d616c"></div><p>Secondary</p></div>
    <div><div class="swatch" style="background:#8b5000"></div><p>Warning</p></div>
    <div><div class="swatch" style="background:#0061a4"></div><p>Info</p></div>
  </div>
  <div class="grid 2" style="margin-top:18px;">
    <div class="card"><h3>單一任務</h3><p>每頁最多一個實心 Primary CTA，引導 Entry → Analysis → Result。</p></div>
    <div class="card"><h3>安全優先</h3><p>空耳警示與「圖不上雲」宣告視為一等 UI，不可藏在設定。</p></div>
  </div>
</section>

<section class="slide">
  <div class="eyebrow">ARCHITECTURE</div>
  <h1>系統架構</h1>
  <img class="figure" src="assets/diagrams/architecture.png" alt="architecture"/>
</section>

<section class="slide">
  <div class="eyebrow">STACK</div>
  <h1>技術棧與工具鏈</h1>
  <div class="grid 3" style="margin-top:24px;">
    <div class="card"><h3>IDE</h3><p>Visual Studio 2026</p></div>
    <div class="card"><h3>App</h3><p>.NET MAUI Android</p></div>
    <div class="card"><h3>API</h3><p>ASP.NET Core Minimal API</p></div>
    <div class="card"><h3>Auth</h3><p>Google＋Email／密碼</p></div>
    <div class="card"><h3>Prompt</h3><p>Google AI Studio</p></div>
    <div class="card"><h3>UI 草稿</h3><p>Google Stitch → Cursor</p></div>
  </div>
</section>

<section class="slide">
  <div class="eyebrow">TIMELINE</div>
  <h1>8 週開發時程</h1>
  <img class="figure" src="assets/diagrams/timeline.png" alt="timeline"/>
</section>

<section class="slide">
  <div class="eyebrow">METRICS & PRIVACY</div>
  <h1>成功指標與隱私</h1>
  <div class="grid" style="grid-template-columns:repeat(5,1fr); margin-top:20px;">
    <div class="card kpi"><div class="v">≥98%</div><p>登入成功率</p></div>
    <div class="card kpi"><div class="v">≥95%</div><p>生成成功率</p></div>
    <div class="card kpi"><div class="v">≥40%</div><p>收藏率</p></div>
    <div class="card kpi"><div class="v">&lt;8s</div><p>文字路徑 P95</p></div>
    <div class="card kpi"><div class="v">≥99%</div><p>Crash-free</p></div>
  </div>
  <div class="grid 2" style="margin-top:18px;">
    <div class="card" style="background:#d0e6f3;"><h3>隱私與資安</h3>
      <ul><li>原圖不離開手機</li><li>僅選定文字送 API／LLM</li><li>LLM Key 不上 App</li><li>JWT／密碼雜湊／最小權限</li></ul>
    </div>
    <div class="card" style="background:#ffdcc0;"><h3>風險對策</h3>
      <ul><li>多語品質不均 → 聲明＋重產＋覆寫語言</li><li>OCR 弱語系 → 手動輸入一等公民</li><li>AI 費用 → 額度＋快取＋熔斷</li><li>金流過早 → Phase 2 再接 Billing</li></ul>
    </div>
  </div>
</section>

<section class="slide">
  <div class="eyebrow">NEXT</div>
  <h1>開工閉環與下一步</h1>
  <div class="big-quote">登入 → 輸入或 OCR 選字 → 多語偵測與空耳 → 選候選 → 存入帳號單字卡</div>
  <div class="grid 3" style="margin-top:22px;">
    <div class="card"><h3>1–2</h3><p>Auth／Users／Quota／JWT</p></div>
    <div class="card"><h3>3–4</h3><p>多語 Prompt＋Analyze＋結果頁</p></div>
    <div class="card"><h3>5–6</h3><p>Notebook＋OCR＋TTS＋送審</p></div>
  </div>
  <p class="sub" style="margin-top:24px;">產品暫名 Soraeru／空耳學單字 · 規劃書 v1.1 · Stitch DESIGN Deep Teal</p>
</section>

</div>
<div class="nav">
  <span>Soraeru｜空耳學單字 MVP</span>
  <span class="hint">← → 或空白鍵切換 · F 全螢幕</span>
  <span id="page">1 / 16</span>
</div>
<script>
  const slides = [...document.querySelectorAll('.slide')];
  let i = 0;
  const page = document.getElementById('page');
  const progress = document.getElementById('progress');
  function show(n) {{
    i = Math.max(0, Math.min(slides.length - 1, n));
    slides.forEach((s, idx) => s.classList.toggle('active', idx === i));
    page.textContent = `${{i + 1}} / ${{slides.length}}`;
    progress.style.width = `${{((i + 1) / slides.length) * 100}}%`;
  }}
  window.addEventListener('keydown', (e) => {{
    if (['ArrowRight','PageDown',' '].includes(e.key)) {{ e.preventDefault(); show(i + 1); }}
    if (['ArrowLeft','PageUp'].includes(e.key)) {{ e.preventDefault(); show(i - 1); }}
    if (e.key === 'Home') show(0);
    if (e.key === 'End') show(slides.length - 1);
    if (e.key.toLowerCase() === 'f') {{
      if (!document.fullscreenElement) document.documentElement.requestFullscreen();
      else document.exitFullscreen();
    }}
  }});
  window.addEventListener('click', (e) => {{
    if (e.clientX > window.innerWidth * 0.55) show(i + 1);
    else show(i - 1);
  }});
  show(0);
</script>
</body>
</html>
"""
    OUT_HTML.write_text(html, encoding="utf-8")
    print(f"HTML -> {OUT_HTML}")


def main():
    print("Building diagrams...")
    make_diagrams()
    print("Building PPTX...")
    build_pptx()
    print("Building HTML...")
    build_html()
    print("All done.")


if __name__ == "__main__":
    main()
